"""
知识库路由 — 纯 HTTP 适配层(业务在 KbService, 存储 SQLite)

迁移后本文件只做请求解析/响应封装与权限上下文注入,
所有业务逻辑(CRUD/解析/检索/产物/成员)在 backend/services/kb_service.py。
API 路径与响应结构与旧版保持兼容, 前端无需改动。

新增接口(权限/共享):
  GET/POST /api/knowledge/notebooks/{kid}/members
  PUT/DELETE /api/knowledge/notebooks/{kid}/members/{user_id}
  POST /api/knowledge/notebooks/{kid}/transfer   — 迁移管理员(唯一)
"""
import json
import re
from fastapi import APIRouter, Depends, UploadFile, File, HTTPException, BackgroundTasks, Request
from fastapi.responses import FileResponse, PlainTextResponse, StreamingResponse
from sqlalchemy.orm import Session
from pydantic import BaseModel

from backend.database import get_db
from backend.core.security import get_optional_user
from backend.models.user import User
from backend.services.kb_service import kb_service
from backend.services.knowledge_parser import extract_text

router = APIRouter()


def _uid(current_user: User | None) -> int:
    """当前用户 id(匿名兜底 0, 由权限层决定是否放行公开库)"""
    return current_user.id if current_user else 0


# ═══════════ 知识库 CRUD ═══════════

@router.get("/notebooks")
def list_kbs(db: Session = Depends(get_db),
             current_user: User = Depends(get_optional_user)):
    """获取可见知识库(我的 + 共享 + 公开)"""
    return kb_service.list_kbs(db, _uid(current_user))


@router.post("/notebooks")
def create_kb(body: dict, db: Session = Depends(get_db),
              current_user: User = Depends(get_optional_user)):
    """创建知识库(创建者自动成为唯一 admin)"""
    return kb_service.create_kb(db, _uid(current_user), body)


@router.put("/notebooks/{kid}")
def update_kb(kid: str, body: dict, db: Session = Depends(get_db),
              current_user: User = Depends(get_optional_user)):
    """更新知识库(editor 及以上)"""
    return kb_service.update_kb(db, kid, _uid(current_user), body)


@router.delete("/notebooks/{kid}")
def delete_kb(kid: str, db: Session = Depends(get_db),
              current_user: User = Depends(get_optional_user)):
    """删除知识库(仅 admin)"""
    return kb_service.delete_kb(db, kid, _uid(current_user))


# ═══════════ 来源 ═══════════

@router.post("/notebooks/{kid}/sources")
async def upload_source(kid: str, background_tasks: BackgroundTasks,
                        file: UploadFile = File(...), db: Session = Depends(get_db),
                        current_user: User = Depends(get_optional_user)):
    """上传来源文档: 格式白名单 + 50MB; 落盘后后台异步解析"""
    content = await file.read()
    return kb_service.upload_source(db, kid, _uid(current_user),
                                    file.filename or "unknown", content, background_tasks)


class TextSourceBody(BaseModel):
    title: str = ""
    content: str = ""


@router.post("/notebooks/{kid}/sources/text")
def add_text_source(kid: str, body: TextSourceBody, db: Session = Depends(get_db),
                    current_user: User = Depends(get_optional_user)):
    """粘贴文本作为来源"""
    return kb_service.add_text_source(db, kid, _uid(current_user), body)


@router.get("/notebooks/{kid}/sources")
def list_sources(kid: str, db: Session = Depends(get_db),
                 current_user: User = Depends(get_optional_user)):
    """获取知识库来源列表"""
    return kb_service.list_sources(db, kid, _uid(current_user))


@router.delete("/notebooks/{kid}/sources/{sid}")
def delete_source(kid: str, sid: str, db: Session = Depends(get_db),
                  current_user: User = Depends(get_optional_user)):
    """删除来源(editor 及以上)"""
    return kb_service.delete_source(db, kid, sid, _uid(current_user))


@router.get("/notebooks/{kid}/sources/{sid}/preview")
def preview_source(kid: str, sid: str, db: Session = Depends(get_db),
                   current_user: User = Depends(get_optional_user)):
    """在线预览来源文件: pdf/html 直接返回文件(iframe 渲染), 文本类返回提取文本"""
    fp, ext, filename = kb_service.get_source_path(db, kid, sid, _uid(current_user))
    if ext in ("pdf", "html", "htm"):
        media = "application/pdf" if ext == "pdf" else "text/html; charset=utf-8"
        return FileResponse(fp, media_type=media)
    if ext in ("txt", "md", "markdown", "csv", "json"):
        try:
            return PlainTextResponse(fp.read_text(encoding="utf-8", errors="replace"))
        except Exception as e:  # noqa: BLE001
            raise HTTPException(400, f"读取失败: {e}")
    if ext == "docx":
        text, err = extract_text(fp, ext)
        if err:
            raise HTTPException(400, err)
        return PlainTextResponse(text)
    raise HTTPException(400, "该格式暂不支持在线预览")


@router.get("/notebooks/{kid}/sources/{sid}/download")
def download_source(kid: str, sid: str, db: Session = Depends(get_db),
                    current_user: User = Depends(get_optional_user)):
    """下载来源原件"""
    fp, _ext, filename = kb_service.get_source_path(db, kid, sid, _uid(current_user))
    return FileResponse(fp, media_type="application/octet-stream", filename=filename)


# ═══════════ 会话 ═══════════

@router.get("/notebooks/{kid}/chats")
def list_chats(kid: str, db: Session = Depends(get_db),
               current_user: User = Depends(get_optional_user)):
    """获取知识库对话历史摘要"""
    return kb_service.list_chats(db, kid, _uid(current_user))


@router.post("/notebooks/{kid}/chats")
def create_chat(kid: str, body: dict = None, db: Session = Depends(get_db),
                current_user: User = Depends(get_optional_user)):
    """创建对话会话"""
    body = body or {}
    return kb_service.create_chat(db, kid, _uid(current_user), body.get("title", ""))


@router.get("/notebooks/{kid}/chats/{cid}")
def get_chat(kid: str, cid: str, db: Session = Depends(get_db),
             current_user: User = Depends(get_optional_user)):
    """获取单个会话完整消息"""
    return kb_service.get_chat(db, kid, cid, _uid(current_user))


@router.delete("/notebooks/{kid}/chats/{cid}")
def delete_chat(kid: str, cid: str, db: Session = Depends(get_db),
                current_user: User = Depends(get_optional_user)):
    """删除对话会话"""
    return kb_service.delete_chat(db, kid, cid, _uid(current_user))


# ═══════════ 问答 ═══════════

class ChatRequest(BaseModel):
    messages: list = []
    sourceIds: list[str] | None = None
    chatId: str | None = None


def _build_context(db: Session, kid: str, source_ids: list[str] | None,
                   limit: int = 5, max_chars: int = 2000) -> tuple[str, list]:
    """构建 RAG/生成上下文: 从 SQLite 读已解析来源文本"""
    from backend.models.kb import KbSource
    sources = db.query(KbSource).filter(
        KbSource.kb_id == kid, KbSource.status == "parsed").all()
    selected = [s for s in sources
                if (not source_ids or s.id in source_ids)][:limit]
    parts = []
    for s in selected:
        text = (s.text_cache or "").strip()
        if not text:
            text = (s.text_preview or "").strip()
        if text:
            parts.append(f"[{s.filename}]\n{text[:max_chars]}")
    context = "\n\n".join(parts) if parts else "（暂无文档内容）"
    return context, [s.to_dict() for s in selected]


def _prepare_chat(db: Session, kid: str, req: ChatRequest, uid: int = 0) -> tuple:
    """构建知识库问答会话 + 压缩历史。返回 (chat, user_msg, history_prefix)。"""
    from backend.core.context_compressor import ContextCompressor
    from backend.models.kb import KbChat

    user_msg = req.messages[-1]["content"] if req.messages else ""

    chat = None
    if req.chatId:
        chat = db.query(KbChat).filter(KbChat.kb_id == kid, KbChat.id == req.chatId).first()
    if chat is None:
        chat = kb_service.create_chat(db, kid, uid, user_msg[:30] or "新对话")
        chat = db.query(KbChat).filter(KbChat.id == chat["id"]).first()

    history_prefix = ""
    hist = [m.to_dict() for m in (chat.messages or [])
            if m.role in ("user", "assistant") and (m.content or "").strip()]
    if hist:
        try:
            compressor = ContextCompressor()
            compressed = compressor.compress_messages(hist[-12:], target_tokens=1500, window_size=8)
            text = "\n".join(f"{m.get('role','')}: {m.get('content','')}"
                             for m in compressed if m.get("content"))
            if text.strip():
                history_prefix = f"## 对话历史(供上下文参考)\n{text[:1200]}"
        except Exception:
            history_prefix = ""
    return chat, user_msg, history_prefix


@router.post("/notebooks/{kid}/chat")
async def chat_kb(kid: str, req: ChatRequest, db: Session = Depends(get_db),
                  current_user: User = Depends(get_optional_user)):
    """知识库问答(分块 RAG, SSE 流式, 会话持久化)。权限: viewer 及以上"""
    kb_service.require_member(db, kid, _uid(current_user), "viewer")
    from backend.services.rag_agent import agentic_chat

    chat, user_msg, history_prefix = _prepare_chat(db, kid, req, _uid(current_user))

    async def event_stream():
        reply = ""
        try:
            async for ev in agentic_chat(kid, req, chat.to_dict(), history_prefix):
                if ev["type"] == "done":
                    reply = ev.get("final_response", "") or reply
                elif ev["type"] == "error":
                    reply = ev.get("message", "") or reply
                elif ev["type"] == "token":
                    reply += ev.get("text", "") or ""
                yield f"event: {ev['type']}\ndata: {json.dumps(ev, ensure_ascii=False)}\n\n"
        except Exception as e:  # noqa: BLE001
            reply = f"（AI 回复暂不可用: {e}）"
            yield f"event: error\ndata: {json.dumps({'type': 'error', 'message': reply}, ensure_ascii=False)}\n\n"
        finally:
            # 持久化必须在 finally 中执行: 客户端断开(GeneratorExit)时同样落盘
            try:
                from backend.models.kb import KbMessage
                chat.messages.append(KbMessage(role="user", content=user_msg))
                chat.messages.append(KbMessage(role="assistant", content=reply or "（无回复）"))
                chat.source_ids = json.dumps(req.sourceIds or [], ensure_ascii=False)
                from datetime import datetime
                chat.updated_at = datetime.now()
                db.commit()
            except Exception:
                pass

    return StreamingResponse(event_stream(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "Connection": "keep-alive"})


# ═══════════ 产物 ═══════════

_GENERATION_EXT = {"html": "html", "mindmap": "md", "ppt": "md", "brief": "md",
                   "timeline": "json"}

GENERATION_LABEL = {
    "html": "网页报告",
    "mindmap": "思维导图",
    "ppt": "PPT 演示",
    "brief": "简报",
    "timeline": "时间轴",
}

MINDMAP_TYPES = ["概括", "时间线", "对比", "结构"]
HTML_STYLES = ["简约白", "商务米色", "科技蓝紫", "深色专业"]
PPT_TEMPLATES = ["简约白", "商务深色", "科技渐变"]


def _build_generation_prompt(kind: str, kb_title: str, context: str,
                             instruction: str, mindmap_type: str,
                             style: str = "", template: str = "") -> str:
    inst = instruction.strip() if instruction else ""
    base = f"你是知识库内容分析专家，仅基于以下资料回答，不要编造。\n\n知识库《{kb_title}》的资料内容：\n{context}\n\n"
    if kind == "html":
        style_hint = f"视觉风格：{style}。" if style in HTML_STYLES else ""
        return base + (
            "请生成一份完整、可直接保存为 .html 的独立研究报告（内联 CSS，含目录、章节、结论）。"
            f"{style_hint}"
            f"标题为《{kb_title}研究报告》。只输出 HTML 代码本身。\n"
            + (f"补充要求：{inst}" if inst else ""))
    if kind == "mindmap":
        types = {"概括": "按主题概括", "时间线": "按时间线", "对比": "按对比维度",
                 "结构": "按层级结构"}
        t = types.get(mindmap_type, "按主题概括")
        return base + (
            f"请用 Markdown 的一级/二级/三级标题层级描述一张思维导图（{t}），"
            f"根节点为《{kb_title}》，分支用 ## / ### 表示。只输出 Markdown。\n"
            + (f"补充要求：{inst}" if inst else ""))
    if kind == "ppt":
        tpl_hint = f"模板风格：{template}。" if template in PPT_TEMPLATES else ""
        return base + (
            "请用 Markdown 写一套演示文稿，每页以一级标题 `# ` 开头作为页标题，"
            f"页内用 ## / ### / - 组织要点，共 8~12 页。{tpl_hint}只输出 Markdown。\n"
            + (f"补充要求：{inst}" if inst else ""))
    if kind == "brief":
        return base + (
            "请输出一份业务简报，Markdown 格式，包含：概述、核心要点(N条)、数据/事实、行动建议、风险。"
            "全文不超过 800 字。只输出 Markdown。\n"
            + (f"补充要求：{inst}" if inst else ""))
    if kind == "timeline":
        return base + (
            "请根据资料中的时间线索，划分不同时间节点，输出 JSON 数组，"
            '每个元素为 {"date": "时间（如 2026-04 或 2026-04-15）", "title": "节点标题", '
            '"detail": "一句话说明"}，按时间升序，5~15 个节点。只输出 JSON 数组本身。\n'
            + (f"补充要求：{inst}" if inst else ""))
    return base


def _build_pptx(md_text: str, template: str, out_path):
    """把产物 markdown(# 开头为页)转成真实 .pptx 文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slides = []
    current = None
    for line in md_text.splitlines():
        if line.startswith("# "):
            current = {"title": line[2:].strip(), "bullets": []}
            slides.append(current)
        elif current is not None and line.strip("#- \t"):
            current["bullets"].append(line.lstrip("#- ").strip())
    if not slides:
        slides = [{"title": "演示文稿", "bullets": [md_text[:200]]}]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    tpl = {
        "商务深色": ((33, 33, 33), (255, 255, 255)),
        "科技渐变": ((24, 32, 68), (235, 240, 255)),
    }.get(template, ((255, 255, 255), (26, 26, 26)))
    bg_rgb, fg_rgb = tpl

    blank = prs.slide_layouts[6]
    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s["title"]
        p.font.size = Pt(36 if i == 0 else 30)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*fg_rgb)
        if s["bullets"]:
            bb = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.8))
            btf = bb.text_frame
            btf.word_wrap = True
            for j, b in enumerate(s["bullets"][:10]):
                p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                p.text = f"• {b}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(*fg_rgb)
                p.space_after = Pt(10)
    prs.save(str(out_path))


@router.post("/notebooks/{kid}/generate")
async def generate_artifact(kid: str, body: dict, db: Session = Depends(get_db),
                            current_user: User = Depends(get_optional_user)):
    """生成知识产出(html/mindmap/ppt/brief/timeline)。权限: editor 及以上"""
    from backend.core.model_config import model_manager
    from backend.config import settings
    from backend.models.kb import KbNotebook, KbArtifact
    import uuid as _uuid
    from pathlib import Path
    from datetime import datetime

    kb_service.require_member(db, kid, _uid(current_user), "editor")

    nb = db.query(KbNotebook).filter(KbNotebook.id == kid).first()
    if not nb:
        raise HTTPException(404, "知识库不存在")

    kind = body.get("kind", "")
    if kind not in _GENERATION_EXT:
        raise HTTPException(400, f"不支持的产出类型: {kind}")

    context, _selected = _build_context(db, kid, body.get("source_ids"),
                                        limit=5, max_chars=8000)
    if context == "（暂无文档内容）":
        return {"degraded": True, "error": "没有可用的已解析来源，请先上传并等待解析完成"}

    prompt = _build_generation_prompt(
        kind, nb.title, context,
        body.get("instruction", ""), body.get("mindmap_type", ""),
        body.get("style", ""), body.get("template", ""),
    )

    try:
        llm = model_manager.get_main_llm()
        content = await llm.chat(
            [{"role": "user", "content": prompt}],
            settings.MAIN_MODEL_NAME, temperature=0.4, max_tokens=4096,
        )
    except Exception as e:  # noqa: BLE001
        return {"degraded": True, "error": f"生成失败: {type(e).__name__}: {e}"}

    content = content.strip()
    content = re.sub(r"^```(?:html|markdown|md|json)?\s*\n?", "", content)
    content = re.sub(r"\n?```\s*$", "", content)

    if kind == "timeline":
        lo, hi = content.find("["), content.rfind("]")
        if lo != -1 and hi > lo:
            candidate = content[lo:hi + 1]
            try:
                json.loads(candidate)
                content = candidate
            except Exception:
                pass

    aid = _uuid.uuid4().hex[:8]
    ext = _GENERATION_EXT[kind]
    filename = f"{aid}_{kind}.{ext}"
    artifacts_dir = _kb_path_artifacts(kid)
    (artifacts_dir / filename).write_text(content, encoding="utf-8")

    art = KbArtifact(
        id=aid, kb_id=kid, kind=kind,
        label=GENERATION_LABEL.get(kind, kind),
        title=nb.title, filename=filename,
        style=body.get("style", ""), template=body.get("template", ""),
        mindmap_type=body.get("mindmap_type", ""),
        prompt=prompt[:200],
    )
    db.add(art)
    db.commit()
    db.refresh(art)

    d = art.to_dict()
    d["url"] = f"/static/knowledge/{kid}/artifacts/{filename}"
    d["markdown"] = content if ext == "md" else ""
    return d


def _kb_path_artifacts(kid: str):
    """产物目录(与 kb_service._artifacts_dir 一致: data/uploads/knowledge/{kid}/artifacts)"""
    from pathlib import Path
    from backend.config import DATA_DIR
    p = DATA_DIR / "uploads" / "knowledge" / kid / "artifacts"
    p.mkdir(parents=True, exist_ok=True)
    return p


@router.get("/notebooks/{kid}/artifacts")
def list_artifacts(kid: str, db: Session = Depends(get_db),
                   current_user: User = Depends(get_optional_user)):
    """获取知识库全部产出"""
    return kb_service.list_artifacts(db, kid, _uid(current_user))


@router.get("/notebooks/{kid}/artifacts/{aid}/pptx")
def download_artifact_pptx(kid: str, aid: str, db: Session = Depends(get_db),
                           current_user: User = Depends(get_optional_user)):
    """把 PPT 产物(markdown)导出为真实 .pptx 下载"""
    from pathlib import Path
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        raise HTTPException(400, "服务端未安装 python-pptx，无法导出 pptx")
    art = kb_service.get_artifact(db, kid, aid, _uid(current_user))
    if art.get("kind") != "ppt":
        raise HTTPException(400, "该产物不是 PPT 类型")
    artifacts_dir = _kb_path_artifacts(kid)
    md_path = artifacts_dir / (art.get("filename") or "")
    if not md_path.exists():
        raise HTTPException(404, "产物文件已丢失")
    out_path = artifacts_dir / f"{aid}.pptx"
    _build_pptx(md_path.read_text(encoding="utf-8"),
                art.get("template", ""), out_path)
    return FileResponse(out_path, media_type="application/octet-stream",
                        filename=f"{art.get('title', '演示文稿')}.pptx")


@router.delete("/notebooks/{kid}/artifacts/{aid}")
def delete_artifact(kid: str, aid: str, db: Session = Depends(get_db),
                    current_user: User = Depends(get_optional_user)):
    """删除产出(editor 及以上)"""
    return kb_service.delete_artifact(db, kid, aid, _uid(current_user))


# ═══════════ 成员 / 共享(新增) ═══════════

@router.get("/notebooks/{kid}/members")
def list_members(kid: str, db: Session = Depends(get_db),
                 current_user: User = Depends(get_optional_user)):
    """成员列表(viewer 及以上可见)"""
    return kb_service.list_members(db, kid, _uid(current_user))


@router.post("/notebooks/{kid}/members")
def add_member(kid: str, body: dict, db: Session = Depends(get_db),
               current_user: User = Depends(get_optional_user)):
    """添加成员(仅 admin); role 限 editor/viewer"""
    return kb_service.add_member(db, kid, _uid(current_user), body)


@router.put("/notebooks/{kid}/members/{user_id}")
def update_member(kid: str, user_id: int, body: dict, db: Session = Depends(get_db),
                  current_user: User = Depends(get_optional_user)):
    """修改成员角色(仅 admin)"""
    return kb_service.update_member_role(db, kid, _uid(current_user), user_id,
                                         body.get("role", ""))


@router.delete("/notebooks/{kid}/members/{user_id}")
def remove_member(kid: str, user_id: int, db: Session = Depends(get_db),
                  current_user: User = Depends(get_optional_user)):
    """移除成员(仅 admin); 不能移除唯一 admin"""
    return kb_service.remove_member(db, kid, _uid(current_user), user_id)


@router.post("/notebooks/{kid}/transfer")
def transfer_admin(kid: str, body: dict, db: Session = Depends(get_db),
                   current_user: User = Depends(get_optional_user)):
    """迁移管理员(仅当前 admin); 迁移后原 admin 降为 editor, 保持唯一"""
    return kb_service.transfer_admin(db, kid, _uid(current_user), body)
