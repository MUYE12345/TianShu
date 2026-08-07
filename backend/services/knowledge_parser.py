"""
知识库来源解析工具 — 从原始文件提取可检索/可预览的纯文本。

支持格式：
  - 文本类：.txt / .md / .csv / .json / .html  → 直接 read_text
  - .pdf    → PyMuPDF(fitz) 逐页提取
  - .docx   → python-docx 提取段落文本
  - .xlsx / .xls → pandas(openpyxl) 逐 sheet 转文本
  - .pptx   → python-pptx 逐页提取文本框
  - 其他格式：返回 ("" , 错误说明)，来源仍保存但不参与检索/预览
"""
from pathlib import Path

PREVIEWABLE_EXTS = {".txt", ".md", ".markdown", ".csv", ".json", ".html", ".htm",
                    ".pdf", ".docx", ".xlsx", ".xls", ".pptx"}


def extract_text(fp: Path, ext: str) -> tuple[str, str | None]:
    """提取文件文本。返回 (text, parse_error)；成功时 parse_error=None。

    解析失败不抛异常，而是返回 ("" 或部分文本, 错误信息)，由调用方标记状态。
    """
    ext = (ext or "").lower().lstrip(".")
    try:
        if ext in ("txt", "md", "markdown", "csv", "json", "html", "htm"):
            return fp.read_text(encoding="utf-8", errors="replace"), None

        if ext == "pdf":
            try:
                import fitz  # PyMuPDF
            except ImportError:
                return "", "未安装 PyMuPDF，无法解析 PDF"
            text = []
            with fitz.open(fp) as doc:
                for page in doc:
                    text.append(page.get_text("text"))
            result = "\n".join(text).strip()
            if not result:
                return "", "PDF 未能提取到文字（可能是扫描件/图片型 PDF）"
            return result, None

        if ext == "docx":
            try:
                from docx import Document
            except ImportError:
                return "", "未安装 python-docx，无法解析 docx"
            doc = Document(str(fp))
            paras = [p.text for p in doc.paragraphs]
            # 表格文字一并提取
            for table in doc.tables:
                for row in table.rows:
                    paras.append(" | ".join(c.text for c in row.cells))
            result = "\n".join(paras).strip()
            if not result:
                return "", "docx 未能提取到文字"
            return result, None

        if ext in ("xlsx", "xls"):
            try:
                import pandas as pd
            except ImportError:
                return "", "未安装 pandas/openpyxl，无法解析 Excel"
            try:
                sheets = pd.read_excel(fp, sheet_name=None, dtype=str)
            except Exception as e:  # noqa: BLE001
                return "", f"Excel 解析失败: {e}"
            parts = []
            for name, df in sheets.items():
                parts.append(f"[工作表: {name}]")
                parts.append(df.fillna("").to_csv(index=False))
            result = "\n".join(parts).strip()
            if not result:
                return "", "Excel 未能提取到内容"
            return result, None

        if ext == "pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return "", "未安装 python-pptx，无法解析 pptx"
            prs = Presentation(str(fp))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text.strip() for sh in slide.shapes
                         if sh.has_text_frame and sh.text.strip()]
                if texts:
                    slides.append(f"[第{i}页]\n" + "\n".join(texts))
            result = "\n\n".join(slides).strip()
            if not result:
                return "", "pptx 未能提取到文字"
            return result, None

        if ext == "doc":
            return "", ".doc 为旧版二进制格式，暂不支持解析，请转换为 .docx 后重新上传"
        if ext == "ppt":
            return "", ".ppt 为旧版二进制格式，暂不支持解析，请转换为 .pptx 后重新上传"

        return "", f"暂不支持解析该格式（.{ext}）"
    except Exception as e:  # noqa: BLE001
        return "", f"解析失败: {type(e).__name__}: {e}"


def is_previewable(ext: str) -> bool:
    """该扩展名是否支持在线预览（提取文本或 iframe 渲染）。"""
    return (ext or "").lower().lstrip(".") in {e.lstrip(".") for e in PREVIEWABLE_EXTS}
